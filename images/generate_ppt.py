"""Generate PPT for 首都师范大学 万爱华 - 农机作业数据处理与分析软件 (~20 slides)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMAGES_DIR = r"D:\PycharmProjects\Agricultural Machinery Operation Data Processing and Analysis Software\images"
OUTPUT_PATH = os.path.join(IMAGES_DIR, "农机作业数据处理与分析软件_万爱华_v2.pptx")

# ====== Unified Color Scheme (Blue-based) ======
PRIMARY      = RGBColor(0x1A, 0x73, 0xE8)   # Main blue
PRIMARY_DARK = RGBColor(0x0D, 0x47, 0x9B)   # Darker blue
PRIMARY_LIGHT= RGBColor(0xE8, 0xF0, 0xFE)   # Very light blue bg
DARK_BG      = RGBColor(0x0A, 0x16, 0x28)   # Dark navy (cover / ending)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT     = RGBColor(0x21, 0x21, 0x21)
MID_TEXT     = RGBColor(0x5C, 0x5C, 0x5C)
LIGHT_BORDER = RGBColor(0xDD, 0xDD, 0xDD)
CARD_BG      = RGBColor(0xF8, 0xF9, 0xFA)
SECTION_BG   = RGBColor(0xF0, 0xF4, 0xF8)
HIGHLIGHT_BG = RGBColor(0xE3, 0xF0, 0xFF)

# ====== Presentation Setup ======
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ====== Helpers ======
def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text; p.font.size = Pt(size); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = font; p.alignment = align
    return tb

def add_multitext(slide, left, top, width, height, lines, size=14, color=DARK_TEXT, spacing=1.4, font='Microsoft YaHei'):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(line, str):
            p.text = line; p.font.size = Pt(size); p.font.color.rgb = color; p.font.name = font
        else:
            p.text = line.get('text',''); p.font.size = Pt(line.get('size', size))
            p.font.color.rgb = line.get('color', color); p.font.bold = line.get('bold', False)
            p.font.name = font
        p.space_after = Pt(size * (spacing - 1))
    return tb

def add_img(slide, name, left, top, width=None, height=None):
    p = os.path.join(IMAGES_DIR, name)
    if os.path.exists(p):
        if width and height: return slide.shapes.add_picture(p, left, top, width, height)
        if width: return slide.shapes.add_picture(p, left, top, width=width)
        if height: return slide.shapes.add_picture(p, left, top, height=height)
    return None

def page_num(slide, n, total):
    add_text(slide, Inches(12.0), Inches(7.05), Inches(1.2), Inches(0.35),
             f'{n}/{total}', size=10, color=MID_TEXT, align=PP_ALIGN.RIGHT)

def header_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.15), PRIMARY)
    add_text(slide, Inches(0.7), Inches(0.2), Inches(11.5), Inches(0.65),
             title, size=26, color=WHITE, bold=True)
    if subtitle:
        add_text(slide, Inches(0.7), Inches(0.7), Inches(11.5), Inches(0.35),
                 subtitle, size=12, color=RGBColor(0xC0, 0xD5, 0xF5))
    # thin accent line
    add_rect(slide, Inches(0), Inches(1.15), Inches(13.333), Inches(0.04), PRIMARY_DARK)

TOTAL = 22  # total slides

# =====================================================================
# SLIDE 1 – 封面
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, DARK_BG)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PRIMARY)
add_text(s, Inches(1.5), Inches(1.8), Inches(10.5), Inches(1.2),
         '农机作业数据处理与分析软件', size=42, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(2.9), Inches(10.5), Inches(0.7),
         'Agricultural Machinery Operation Data Processing and Analysis Software',
         size=18, color=RGBColor(0x90,0xAA,0xCC), align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.8), Inches(3.8), Inches(1.733), Inches(0.04), PRIMARY)
add_text(s, Inches(1.5), Inches(4.2), Inches(10.5), Inches(0.6),
         '首都师范大学', size=24, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(4.8), Inches(10.5), Inches(0.6),
         '万爱华', size=20, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
add_rect(s, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), PRIMARY)
page_num(s, 1, TOTAL)

# =====================================================================
# SLIDE 2 – 目录
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '目  录', 'CONTENTS')
toc = [
    ('01', '项目背景与需求'),('02', '数据说明'),('03', '系统架构设计'),
    ('04', '技术选型'),('05', '数据库设计（一）'),('06', '数据库设计（二）'),
    ('07', '功能模块总览'),('08', '用户认证与安全'),('09', '轨迹列表与管理'),
    ('10', '地图展示 — 卫星与矢量底图'),('11', '地图交互 — 动态回放与测距'),
    ('12', '数据导入'),('13', '文件管理'),('14', '用户管理'),
    ('15', '角色与权限管理（RBAC）'),('16', '菜单权限配置'),
    ('17', '技术亮点（一）'),('18', '技术亮点（二）'),
    ('19', '项目总结'),('20', '感谢聆听'),
]
for i, (num, title) in enumerate(toc):
    col = i % 2; row = i // 2
    x = Inches(1.2) + Inches(col * 6.0)
    y = Inches(1.55) + Inches(row * 0.55)
    shape = s.shapes.add_shape(MSO_SHAPE.OVAL, x, y+Inches(0.03), Inches(0.32), Inches(0.32))
    shape.fill.solid(); shape.fill.fore_color.rgb = PRIMARY; shape.line.fill.background()
    tf = shape.text_frame; tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(10); tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, x+Inches(0.45), y, Inches(4.5), Inches(0.35), title, size=14, color=DARK_TEXT)
page_num(s, 2, TOTAL)

# =====================================================================
# SLIDE 3 – 项目背景与需求
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '01  项目背景与需求', 'BACKGROUND & REQUIREMENTS')
add_text(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.45), '实践内容', size=20, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.5), [
    '农机作业数据处理与分析软件的基础功能设计与开发',
    '',
    '• 设计并实现一个完整的Web应用程序',
    '• 采用数据库存储原始农机轨迹数据',
    '• 实现用户和管理员两种角色',
    '• 集成地图组件展示轨迹数据',
    '• 独立完成需求分析、设计、编码、验收全流程',
], size=13, color=MID_TEXT)
add_text(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.45), '功能需求总览', size=20, color=PRIMARY, bold=True)
add_multitext(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(1.5), [
    {'text':'用户基本功能','size':14,'bold':True,'color':DARK_TEXT},
    '• 注册、登录、修改密码',
    '• 地图操作：放大、缩小、漫游、测距',
    '• 加载遥感影像与矢量地图',
    '• 轨迹数据可视化与属性查询',
    '',
    {'text':'管理员功能','size':14,'bold':True,'color':DARK_TEXT},
    '• Excel轨迹数据导入   • 轨迹数据删除',
    '• 用户账号管理        • 角色与菜单权限管理',
], size=13, color=MID_TEXT)
# bottom note
add_rect(s, Inches(0.7), Inches(6.2), Inches(11.9), Inches(0.65), PRIMARY_LIGHT)
add_text(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.45),
         '编程要求：采用Web应用程序开发，前后端分离架构，使用数据库存储原始轨迹数据，集成地图组件',
         size=12, color=PRIMARY_DARK)
page_num(s, 3, TOTAL)

# =====================================================================
# SLIDE 4 – 数据说明
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '02  数据说明', 'DATA DESCRIPTION')
add_text(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.45), '原始数据', size=20, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.7), Inches(2.1), Inches(5.5), Inches(2.0), [
    '• 提供8个Excel格式的农机地块作业轨迹文件',
    '• 每个文件对应一条完整的作业轨迹记录',
    '• 单条轨迹包含3000-8000+个GPS采样点',
    '• 数据采集频率为1Hz（每秒1个点）',
], size=13, color=MID_TEXT)
add_text(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.45), '数据字段', size=20, color=PRIMARY, bold=True)
add_multitext(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.0), [
    {'text':'GPS定位数据','size':14,'bold':True,'color':DARK_TEXT},
    '    gpstime    GPS时间戳           lon/lat    经纬度坐标',
    '',
    {'text':'作业状态数据','size':14,'bold':True,'color':DARK_TEXT},
    '    workstatus  作业状态（1=正常作业, 0=闲置停止）',
    '    depth       耕深（cm）         depthstandard  标准耕深',
    '',
    {'text':'运动数据','size':14,'bold':True,'color':DARK_TEXT},
    '    velocity    瞬时速度（km/h）    course    行进方向角',
    '    width       作业幅宽（m）',
], size=12, color=MID_TEXT)
# data flow
add_rect(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.0), PRIMARY_LIGHT)
add_text(s, Inches(0.9), Inches(5.55), Inches(11.5), Inches(0.35),
         '数据流程', size=16, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.9), Inches(5.95), Inches(11.5), Inches(0.45), [
    'Excel文件  ──上传──▶  后端解析  ──存储──▶  track表 + trackpoints表  ──计算──▶  work表 + rate表  ──API──▶  前端地图展示',
], size=13, color=PRIMARY_DARK)
page_num(s, 4, TOTAL)

# =====================================================================
# SLIDE 5 – 系统架构设计
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '03  系统架构设计', 'SYSTEM ARCHITECTURE')
# three-tier architecture
# Frontend
add_rounded_rect(s, Inches(1.0), Inches(1.8), Inches(4.8), Inches(2.2), PRIMARY_LIGHT)
add_rect(s, Inches(1.0), Inches(1.8), Inches(4.8), Inches(0.5), PRIMARY)
add_text(s, Inches(1.2), Inches(1.88), Inches(4.4), Inches(0.35), '前端展示层  Frontend', size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_multitext(s, Inches(1.3), Inches(2.55), Inches(4.2), Inches(1.3), [
    '框架：Vue 3 + Vite',
    'UI库：Element Plus',
    '地图：Leaflet + 高德瓦片',
    '路由：Vue Router',
    'HTTP：Axios',
], size=13, color=DARK_TEXT)
# Backend
add_rounded_rect(s, Inches(7.3), Inches(1.8), Inches(4.8), Inches(2.2), HIGHLIGHT_BG)
add_rect(s, Inches(7.3), Inches(1.8), Inches(4.8), Inches(0.5), PRIMARY_DARK)
add_text(s, Inches(7.5), Inches(1.88), Inches(4.4), Inches(0.35), '后端服务层  Backend', size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_multitext(s, Inches(7.6), Inches(2.55), Inches(4.2), Inches(1.3), [
    '框架：Django 4.2 + DRF',
    'API设计：RESTful',
    '数据库：SQLite + Django ORM',
    '认证：Token / Session',
    '跨域：django-cors-headers',
], size=13, color=DARK_TEXT)
# Middle arrow
add_text(s, Inches(5.6), Inches(2.5), Inches(2.0), Inches(0.6), '◄  REST API  ►', size=13, color=MID_TEXT, align=PP_ALIGN.CENTER)
# Database layer
add_rounded_rect(s, Inches(2.5), Inches(4.7), Inches(8.3), Inches(1.2), CARD_BG)
add_text(s, Inches(2.7), Inches(4.8), Inches(7.8), Inches(0.35), '数据层  Database', size=18, color=PRIMARY_DARK, bold=True)
add_multitext(s, Inches(2.7), Inches(5.25), Inches(7.8), Inches(0.55), [
    'SQLite 数据库 | 10张数据表 | Django ORM 对象关系映射 | 支持批量导入与事务管理',
], size=13, color=MID_TEXT)
# bottom note
add_rect(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.55), SECTION_BG)
add_text(s, Inches(0.9), Inches(6.42), Inches(11.5), Inches(0.4),
         '架构优势：前后端完全解耦 ▶ 独立开发部署 ▶ 前端响应式UI ▶ 后端RESTful接口 ▶ 轻量级数据库',
         size=12, color=PRIMARY_DARK)
page_num(s, 5, TOTAL)

# =====================================================================
# SLIDE 6 – 技术选型
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '04  技术选型', 'TECHNOLOGY STACK')
techs = [
    ('Python 3', '后端开发语言，Django框架基础', '语法简洁，生态丰富'),
    ('Django 4.2', 'Web应用框架，ORM + DRF + Admin', '快速开发，安全可靠'),
    ('DRF', 'Django REST Framework，构建API', '序列化、视图集、路由'),
    ('Vue 3', '前端渐进式框架，Composition API', '响应式，组件化开发'),
    ('Vite', '新一代前端构建工具', '极速HMR，按需编译'),
    ('Element Plus', 'Vue 3 UI组件库', '丰富组件，统一风格'),
    ('Leaflet', '开源地图库', '轻量，插件丰富'),
    ('SQLite', '嵌入式关系型数据库', '零配置，单文件存储'),
]
for i, (name, desc, why) in enumerate(techs):
    col = i % 4; row = i // 4
    x = Inches(0.5) + Inches(col * 3.15); y = Inches(1.55) + Inches(row * 2.8)
    add_rounded_rect(s, x, y, Inches(2.9), Inches(2.5), CARD_BG)
    add_rect(s, x, y, Inches(2.9), Inches(0.5), PRIMARY)
    add_text(s, x+Inches(0.1), y+Inches(0.06), Inches(2.7), Inches(0.38),
             name, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_multitext(s, x+Inches(0.15), y+Inches(0.65), Inches(2.6), Inches(1.7), [
        {'text':desc,'size':11,'color':DARK_TEXT},
        '',
        {'text':'选型理由：'+why,'size':10,'color':MID_TEXT},
    ], size=11, color=DARK_TEXT)
page_num(s, 6, TOTAL)

# =====================================================================
# SLIDE 7 – 数据库设计（一）
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '05  数据库设计（一）— 核心业务表', 'DATABASE DESIGN – CORE TABLES')
tables = [
    ('track', '轨迹主表', 'trackid(PK), starttime, endtime, width, totalpoints', '每条轨迹一条记录'),
    ('trackpoints', '轨迹点表', 'id(PK), trackid(FK), gpstime, lon, lat, velocity, depth, workstatus, width', '每个GPS点一条记录，1:N'),
    ('work', '作业统计表', 'trackid(PK+FK), worktime, worklength, workarea, avgvelocity', '作业时长/面积/速度统计，1:1'),
    ('rate', '通行率表', 'trackid(PK+FK), passrate, productionrate, timerrate', '通行率/生产效率指标，1:1'),
    ('user', '用户表', 'id(PK), username, password, nickname, email, phone, role(FK)', '用户基本信息'),
    ('role', '角色表', 'id(PK), name, description, flag', '角色定义（管理员/用户）'),
    ('menu', '菜单表', 'id(PK), name, path, icon, pid(FK自引用), page_path', '树形菜单结构'),
    ('role_menu', '角色菜单关联', 'role_id(PK+FK), menu_id(FK)', '角色与菜单多对多'),
    ('file', '文件表', 'id(PK), name, type, size, url, md5', '上传文件管理'),
    ('import_log', '导入日志表', 'id(PK), admin_id(FK), file_name, import_count, import_status, import_time', '数据导入日志'),
]
for i, (name, desc, fields, note) in enumerate(tables):
    y = Inches(1.45) + Inches(i * 0.55)
    color = PRIMARY if i < 4 else PRIMARY_DARK
    badge = add_rounded_rect(s, Inches(0.5), y, Inches(1.8), Inches(0.35), color)
    tf = badge.text_frame; tf.paragraphs[0].text = name
    tf.paragraphs[0].font.size = Pt(11); tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True; tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_text(s, Inches(2.45), y, Inches(1.2), Inches(0.35), desc, size=11, color=DARK_TEXT, bold=True)
    add_text(s, Inches(3.8), y, Inches(6.5), Inches(0.35), fields, size=9, color=MID_TEXT)
    add_text(s, Inches(10.5), y, Inches(2.3), Inches(0.35), note, size=9, color=PRIMARY_DARK)
page_num(s, 7, TOTAL)

# =====================================================================
# SLIDE 8 – 数据库设计（二）ER关系
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '06  数据库设计（二）— 关系与约束', 'DATABASE DESIGN – RELATIONSHIPS')
# ER diagram area
add_rect(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(3.2), CARD_BG)
add_text(s, Inches(0.7), Inches(1.7), Inches(5.0), Inches(0.35), '表关系图（ER）', size=16, color=PRIMARY, bold=True)
# Boxes for ER
er_boxes = [
    ('track', Inches(0.8), Inches(2.3)),
    ('trackpoints', Inches(3.5), Inches(2.3)),
    ('work', Inches(0.8), Inches(3.6)),
    ('rate', Inches(3.5), Inches(3.6)),
    ('user', Inches(7.0), Inches(2.3)),
    ('role', Inches(7.0), Inches(3.6)),
    ('menu', Inches(9.5), Inches(2.3)),
    ('role_menu', Inches(9.5), Inches(3.6)),
    ('file', Inches(0.8), Inches(4.5)),
    ('import_log', Inches(3.5), Inches(4.5)),
]
for name, x, y in er_boxes:
    bx = add_rounded_rect(s, x, y, Inches(2.0), Inches(0.55), WHITE)
    bx.line.color.rgb = PRIMARY; bx.line.width = Pt(1)
    add_text(s, x+Inches(0.1), y+Inches(0.1), Inches(1.8), Inches(0.35), name, size=13, color=PRIMARY, bold=True, align=PP_ALIGN.CENTER)

add_text(s, Inches(0.7), Inches(5.15), Inches(12.0), Inches(0.35), '关键关系与约束', size=16, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.7), Inches(5.55), Inches(5.5), Inches(1.3), [
    {'text':'一对多 (1:N)','size':13,'bold':True,'color':DARK_TEXT},
    'track → trackpoints：一条轨迹包含多个轨迹点，通过trackid外键关联',
    'role → user：一个角色可分配给多个用户',
    {'text':'一对一 (1:1)','size':13,'bold':True,'color':DARK_TEXT},
    'track → work：每条轨迹对应一个作业统计记录（主键即为外键）',
    'track → rate：每条轨迹对应一个通行率记录',
], size=11, color=MID_TEXT)
add_multitext(s, Inches(7.0), Inches(5.55), Inches(5.8), Inches(1.3), [
    {'text':'多对多 (N:M)','size':13,'bold':True,'color':DARK_TEXT},
    'role ↔ menu：通过role_menu中间表实现角色与菜单的多对多关联',
    {'text':'自引用','size':13,'bold':True,'color':DARK_TEXT},
    'menu.pid → menu.id：菜单表通过pid字段实现父子层级关系，支持多级菜单树',
    {'text':'级联删除','size':13,'bold':True,'color':DARK_TEXT},
    '删除track时级联删除对应的trackpoints、work、rate记录',
], size=11, color=MID_TEXT)
page_num(s, 8, TOTAL)

# =====================================================================
# SLIDE 9 – 功能模块总览
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '07  功能模块总览', 'FEATURE MODULES OVERVIEW')
modules = [
    ('用户认证', '注册 / 登录\n修改密码 / 个人信息\nToken认证', PRIMARY),
    ('轨迹管理', '轨迹列表 / 筛选\n查看轨迹详情\n管理员删除轨迹', PRIMARY_DARK),
    ('地图展示', '卫星影像底图\n矢量地图底图\n轨迹叠加着色', PRIMARY),
    ('地图交互', '放大缩小漫游\n距离测量工具\n轨迹动态回放', PRIMARY_DARK),
    ('数据导入', 'Excel上传解析\n批量数据导入\n导入日志记录', PRIMARY),
    ('权限管理', '用户管理\n角色管理\n菜单权限配置', PRIMARY_DARK),
]
for i, (title, desc, color) in enumerate(modules):
    col = i % 3; row = i // 3
    x = Inches(0.6) + Inches(col * 4.15); y = Inches(1.55) + Inches(row * 2.85)
    add_rounded_rect(s, x, y, Inches(3.85), Inches(2.55), CARD_BG)
    add_rect(s, x, y, Inches(3.85), Inches(0.55), color)
    add_text(s, x+Inches(0.15), y+Inches(0.08), Inches(3.55), Inches(0.4),
             title, size=17, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x+Inches(0.2), y+Inches(0.8), Inches(3.45), Inches(1.6),
             desc, size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)
page_num(s, 9, TOTAL)

# =====================================================================
# SLIDE 10 – 用户认证与安全
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '08  用户认证与安全', 'AUTHENTICATION & SECURITY')
add_img(s, '个人信息修改页面.png', Inches(0.5), Inches(1.5), width=Inches(5.8))
add_img(s, '修改密码界面.png', Inches(6.8), Inches(1.5), width=Inches(5.8))
add_text(s, Inches(0.5), Inches(5.35), Inches(5.8), Inches(0.3),
         '▲ 个人信息修改 — 昵称、邮箱、电话、地址', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.8), Inches(5.35), Inches(5.8), Inches(0.3),
         '▲ 修改密码 — 旧密码验证 + 新密码确认', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.35), '认证机制', size=17, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.8), [
    '• 用户登录后前端存储Token，Axios请求拦截器自动附加Authorization头    • 路由守卫（beforeEach）拦截未登录访问，自动跳转登录页',
    '• 密码采用加密存储，修改密码需验证旧密码    • 注册时进行用户名唯一性校验    • 个人信息修改包含昵称、邮箱、电话、地址等字段',
], size=11, color=MID_TEXT)
page_num(s, 10, TOTAL)

# =====================================================================
# SLIDE 11 – 轨迹列表与管理
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '09  轨迹列表与管理', 'TRACK LIST & MANAGEMENT')
add_img(s, '轨迹管理页面.png', Inches(0.3), Inches(1.5), width=Inches(7.8))
add_text(s, Inches(0.3), Inches(5.5), Inches(7.8), Inches(0.3),
         '▲ 轨迹列表页 — 支持时间范围、幅宽筛选，分页浏览', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(0.4), '功能要点', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(8.5), Inches(2.1), Inches(4.3), Inches(3.5), [
    {'text':'数据列表','size':14,'bold':True,'color':DARK_TEXT},
    '展示全部轨迹记录（轨迹ID、起止时间、幅宽、总点数），支持分页浏览',
    '',
    {'text':'多维筛选','size':14,'bold':True,'color':DARK_TEXT},
    '• 按时间范围筛选（起止日期）',
    '• 按幅宽范围筛选（最小/最大）',
    '• 筛选条件可组合使用',
    '',
    {'text':'操作功能','size':14,'bold':True,'color':DARK_TEXT},
    '• 查看轨迹：跳转地图页展示轨迹',
    '• 删除轨迹：管理员权限，级联删除关联数据',
    '• 上传轨迹：跳转数据导入页',
], size=11, color=MID_TEXT)
page_num(s, 11, TOTAL)

# =====================================================================
# SLIDE 12 – 地图展示（卫星与矢量底图）
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '10  地图展示 — 卫星与矢量底图', 'MAP DISPLAY – SATELLITE & VECTOR BASEMAP')
add_img(s, '轨迹展示页面.png', Inches(0.3), Inches(1.5), width=Inches(6.3))
add_img(s, '轨迹矢量展示页面.png', Inches(6.9), Inches(1.5), width=Inches(6.1))
add_text(s, Inches(0.3), Inches(5.7), Inches(6.3), Inches(0.3),
         '▲ 卫星底图模式 — 遥感影像叠加轨迹矢量', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.9), Inches(5.7), Inches(6.1), Inches(0.3),
         '▲ 矢量底图模式 — 高德标准地图叠加轨迹', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.3), Inches(6.15), Inches(12.5), Inches(0.35), '底图切换与轨迹显示', size=17, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.3), Inches(6.55), Inches(12.5), Inches(0.7), [
    '• 左上角"图层"按钮一键切换卫星影像 / 矢量地图    • 轨迹按作业状态着色：绿色=正常作业（workstatus=1），灰色=闲置停止（workstatus=0）',
    '• WGS-84 → GCJ-02坐标自动转换，确保轨迹与高德底图精准叠加    • 点击轨迹点位查看详情面板（经纬度、速度、耕深）    • 右侧统计卡片显示作业时长、面积、平均速度',
], size=11, color=MID_TEXT)
page_num(s, 12, TOTAL)

# =====================================================================
# SLIDE 13 – 动态回放与测距
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '11  地图交互 — 动态回放与测距', 'MAP INTERACTION – PLAYBACK & MEASUREMENT')
add_text(s, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.4), '轨迹动态回放', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.5), Inches(2.05), Inches(5.8), Inches(2.5), [
    '• 点击"动态运行展示"按钮启动轨迹回放',
    '• 按GPS时间序列逐点渲染，每条轨迹最长播放数百秒',
    '• 使用setInterval精确控制播放速度（0.1秒/点）',
    '• 当前点以高亮标记（金色圆点+白色边框）',
    '• 已播放点形成金色轨迹线，实现"拖尾"效果',
    '• 地图自动跟随（panTo平滑动画）',
    '• 点击"停止展示"可随时中断回放',
    '',
    '技术要点：基于Leaflet的LayerGroup动态管理回放图层，',
    '每次interval触发时clearLayers后重新绘制，确保性能。',
], size=11, color=MID_TEXT)
add_text(s, Inches(7.0), Inches(1.5), Inches(6.0), Inches(0.4), '距离测量', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(7.0), Inches(2.05), Inches(6.0), Inches(2.5), [
    '• 点击"测距"按钮进入测量模式',
    '• 左键在地图上添加测量点，右键撤回上一个点',
    '• 实时计算并显示累计距离',
    '• 自动单位转换：<1000m显示"xxx.x m"，>=1000m显示"x.xxx km"',
    '• 测量线以橙色虚线显示，每个测量点带编号标记',
    '• 支持连续添加任意多个测量点',
    '• 点击"结束测距"退出，"清空测距"清除所有标记',
    '',
    '实现方式：利用Leaflet的L.latLng.distanceTo()方法',
    '计算相邻点间球面距离，精度可靠。',
], size=11, color=MID_TEXT)
add_rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.8), PRIMARY_LIGHT)
add_multitext(s, Inches(0.7), Inches(6.08), Inches(11.9), Inches(0.65), [
    {'text':'交互设计原则','size':13,'bold':True,'color':PRIMARY_DARK},
    '回放与测距互斥（进入测距时自动停止回放）| 测距期间隐藏轨迹默认渲染，突出测量内容 | 所有图层操作均不干扰底图切换功能',
], size=11, color=MID_TEXT)
page_num(s, 13, TOTAL)

# =====================================================================
# SLIDE 14 – 数据导入
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '12  数据导入', 'DATA IMPORT')
add_img(s, '上传日志页面.png', Inches(0.3), Inches(1.5), width=Inches(7.8))
add_text(s, Inches(0.3), Inches(5.55), Inches(7.8), Inches(0.3),
         '▲ 数据导入页 — 上传Excel + 实时日志反馈', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(0.4), '导入流程', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(8.5), Inches(2.1), Inches(4.3), Inches(4.0), [
    {'text':'Step 1  文件上传','size':14,'bold':True,'color':DARK_TEXT},
    '管理员选择Excel文件，前端校验文件格式（.xls/.xlsx）',
    '',
    {'text':'Step 2  数据解析','size':14,'bold':True,'color':DARK_TEXT},
    '后端使用openpyxl/pandas读取Excel，逐行提取GPS数据，验证经纬度、时间等关键字段',
    '',
    {'text':'Step 3  数据入库','size':14,'bold':True,'color':DARK_TEXT},
    '创建track记录 ▶ 批量插入trackpoints ▶ 计算work统计 ▶ 计算rate指标',
    '',
    {'text':'Step 4  日志记录','size':14,'bold':True,'color':DARK_TEXT},
    '生成ImportLog记录（文件名、导入数量、状态、错误信息），前端实时展示',
], size=10, color=MID_TEXT)
page_num(s, 14, TOTAL)

# =====================================================================
# SLIDE 15 – 文件管理
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '13  文件管理', 'FILE MANAGEMENT')
add_img(s, '文件管理页面.png', Inches(0.3), Inches(1.5), width=Inches(7.8))
add_text(s, Inches(0.3), Inches(5.35), Inches(7.8), Inches(0.3),
         '▲ 文件管理页 — 文件列表、上传、删除', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(0.4), '功能要点', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(8.5), Inches(2.1), Inches(4.3), Inches(3.5), [
    {'text':'文件上传','size':14,'bold':True,'color':DARK_TEXT},
    '• 支持通用文件上传',
    '• 记录文件类型、大小、MD5校验',
    '',
    {'text':'文件列表','size':14,'bold':True,'color':DARK_TEXT},
    '• 展示所有已上传文件',
    '• 显示文件名、类型、大小、URL',
    '',
    {'text':'文件操作','size':14,'bold':True,'color':DARK_TEXT},
    '• 单文件删除 / 批量删除',
    '• 文件MD5去重校验',
    '',
    {'text':'与数据导入的关系','size':14,'bold':True,'color':DARK_TEXT},
    '文件管理是通用的文件存储模块；',
    '数据导入是专门的Excel轨迹导入；',
    '两者协同工作，实现完整的数据管理链路。',
], size=10, color=MID_TEXT)
page_num(s, 15, TOTAL)

# =====================================================================
# SLIDE 16 – 用户管理
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '14  用户管理', 'USER MANAGEMENT')
add_img(s, '用户管理页面.png', Inches(0.3), Inches(1.5), width=Inches(7.8))
add_text(s, Inches(0.3), Inches(5.6), Inches(7.8), Inches(0.3),
         '▲ 用户管理页 — 查看所有用户、按角色筛选', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(8.5), Inches(1.5), Inches(4.3), Inches(0.4), '管理功能', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(8.5), Inches(2.1), Inches(4.3), Inches(4.0), [
    {'text':'用户列表','size':14,'bold':True,'color':DARK_TEXT},
    '• 展示所有注册用户信息',
    '• 包含用户名、昵称、邮箱、电话、角色、注册时间',
    '• 支持分页浏览',
    '',
    {'text':'筛选与搜索','size':14,'bold':True,'color':DARK_TEXT},
    '• 按角色类型筛选用户',
    '• 支持关键词搜索',
    '',
    {'text':'操作功能','size':14,'bold':True,'color':DARK_TEXT},
    '• 编辑用户信息',
    '• 删除用户（支持批量删除）',
    '• 管理员可管理所有用户',
], size=11, color=MID_TEXT)
page_num(s, 16, TOTAL)

# =====================================================================
# SLIDE 17 – 角色与权限管理（RBAC）
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '15  角色与权限管理（RBAC）', 'ROLE-BASED ACCESS CONTROL')
add_img(s, '角色管理页面.png', Inches(0.3), Inches(1.5), width=Inches(6.0))
add_text(s, Inches(0.3), Inches(4.9), Inches(6.0), Inches(0.3),
         '▲ 角色管理 — 创建、编辑、删除角色', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.4), 'RBAC模型设计', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(3.5), [
    {'text':'核心概念','size':14,'bold':True,'color':DARK_TEXT},
    '• 用户（User）→ 角色（Role）→ 菜单权限（Menu）',
    '• 用户与角色：多对一（一个用户属于一个角色）',
    '• 角色与菜单：多对多（通过role_menu中间表）',
    '',
    {'text':'权限控制流程','size':14,'bold':True,'color':DARK_TEXT},
    '① 用户登录 → 获取所属角色',
    '② 根据角色 → 查询role_menu → 获取菜单权限列表',
    '③ 前端路由守卫 → 比对当前路径与权限 → 放行/拦截',
    '④ 后端API → 响应头携带X-User-Role → 业务逻辑鉴权',
    '',
    {'text':'实现效果','size':14,'bold':True,'color':DARK_TEXT},
    '不同角色登录后，侧边栏菜单动态渲染（仅显示有权限的菜单项）',
], size=10, color=MID_TEXT)
page_num(s, 17, TOTAL)

# =====================================================================
# SLIDE 18 – 菜单权限配置
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '16  菜单权限配置', 'MENU PERMISSION CONFIGURATION')
add_img(s, '菜单管理页面.png', Inches(0.3), Inches(1.5), width=Inches(6.0))
add_text(s, Inches(0.3), Inches(5.0), Inches(6.0), Inches(0.3),
         '▲ 菜单管理 — 树形菜单结构 + 图标选择', size=10, color=MID_TEXT, align=PP_ALIGN.CENTER)
add_text(s, Inches(6.8), Inches(1.5), Inches(6.0), Inches(0.4), '菜单设计', size=19, color=PRIMARY, bold=True)
add_multitext(s, Inches(6.8), Inches(2.1), Inches(6.0), Inches(4.0), [
    {'text':'树形结构','size':14,'bold':True,'color':DARK_TEXT},
    '• menu表通过pid字段实现自引用父子关系',
    '• 支持多级菜单嵌套（如：轨迹管理 > 轨迹列表）',
    '• 每个菜单项定义：名称、路径、图标、页面组件',
    '',
    {'text':'权限分配','size':14,'bold':True,'color':DARK_TEXT},
    '• 在角色管理页面为每个角色勾选可访问的菜单',
    '• 数据存储在role_menu关联表',
    '',
    {'text':'动态渲染','size':14,'bold':True,'color':DARK_TEXT},
    '• 前端根据当前用户角色查询菜单树',
    '• Element Plus el-menu组件动态渲染侧边栏',
    '• 无权限的页面路由被守卫拦截',
], size=10, color=MID_TEXT)
page_num(s, 18, TOTAL)

# =====================================================================
# SLIDE 19 – 技术亮点（一）
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '17  技术亮点（一）', 'TECHNICAL HIGHLIGHTS I')
hl = [
    ('WGS-84 → GCJ-02\n坐标转换算法',
     '中国大陆地图使用GCJ-02坐标系，GPS原始数据为WGS-84坐标系。项目实现了完整的坐标转换算法，包含20+个三角函数运算，将GPS坐标精确转换为高德地图坐标系，偏移精度控制在米级，确保轨迹与底图完美叠加。'),
    ('动态轨迹回放引擎',
     '基于setInterval实现逐点回放引擎，支持启动/停止控制。回放过程中地图自动跟随（panTo平滑动画），已走过路径以金色轨迹线标记，当前点以高亮标记突出显示，形成流畅的"拖尾"可视化效果。'),
    ('底图实时切换',
     '通过自定义Leaflet Control实现卫星影像与矢量地图的一键切换。使用Layer Management API动态替换底图TileLayer，无需刷新页面。配合CSS动画实现按钮状态平滑过渡，交互体验流畅。'),
]
for i, (title, desc) in enumerate(hl):
    y = Inches(1.5) + Inches(i * 1.9)
    add_rounded_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.7), CARD_BG)
    add_rect(s, Inches(0.5), y, Inches(0.06), Inches(1.7), PRIMARY)
    add_text(s, Inches(0.8), y+Inches(0.1), Inches(3.0), Inches(0.8), title, size=15, color=PRIMARY, bold=True)
    add_text(s, Inches(4.0), y+Inches(0.15), Inches(8.5), Inches(1.4), desc, size=12, color=MID_TEXT)
page_num(s, 19, TOTAL)

# =====================================================================
# SLIDE 20 – 技术亮点（二）
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '18  技术亮点（二）', 'TECHNICAL HIGHLIGHTS II')
hl2 = [
    ('RBAC权限控制系统',
     '实现完整的用户-角色-菜单三级权限模型。前端路由守卫（beforeEach）+ 后端API双重鉴权。侧边栏菜单根据用户角色动态渲染，无权限页面自动拦截。采用JWT Token机制，Axios拦截器自动附加认证信息。'),
    ('大数据量SVG渲染优化',
     '单条轨迹包含4000+个GPS点，全部以SVG path元素渲染到地图上。通过按workstatus分组（分段渲染），减少DOM操作次数。使用Leaflet SVG渲染器替代Canvas，确保轨迹线清晰可见且交互响应流畅。'),
    ('前后端分离架构',
     'Vue3 + Vite前端与Django REST Framework后端完全解耦，通过HTTP/REST API通信。前端使用Composition API组织代码，后端使用ViewSet + Router自动生成RESTful接口。CORS跨域配置完善，开发与部署独立进行。'),
]
for i, (title, desc) in enumerate(hl2):
    y = Inches(1.5) + Inches(i * 1.9)
    add_rounded_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.7), CARD_BG)
    add_rect(s, Inches(0.5), y, Inches(0.06), Inches(1.7), PRIMARY_DARK)
    add_text(s, Inches(0.8), y+Inches(0.1), Inches(3.0), Inches(0.8), title, size=15, color=PRIMARY, bold=True)
    add_text(s, Inches(4.0), y+Inches(0.15), Inches(8.5), Inches(1.4), desc, size=12, color=MID_TEXT)
page_num(s, 20, TOTAL)

# =====================================================================
# SLIDE 21 – 项目总结
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, WHITE)
header_bar(s, '19  项目总结', 'PROJECT SUMMARY')
add_text(s, Inches(0.7), Inches(1.5), Inches(5.5), Inches(0.45), '已完成的工作', size=20, color=PRIMARY, bold=True)
add_multitext(s, Inches(0.7), Inches(2.1), Inches(5.5), Inches(3.8), [
    '✓ 需求分析：完成功能需求、数据需求、角色需求的全面分析',
    '✓ 数据库设计：设计10张规范化数据表，涵盖轨迹、用户、权限',
    '✓ 系统架构：Django + Vue3前后端分离架构',
    '✓ 用户认证：注册/登录/修改密码/个人信息，JWT Token机制',
    '✓ 轨迹管理：8条轨迹数据成功导入，列表筛选与分页展示',
    '✓ 地图展示：卫星影像+矢量地图双底图，轨迹按作业状态着色',
    '✓ 地图交互：放大缩小漫游、距离测量、轨迹动态回放',
    '✓ 坐标转换：WGS-84 → GCJ-02精确坐标转换算法',
    '✓ 数据导入：Excel文件上传解析，批量入库+日志记录',
    '✓ RBAC权限：用户-角色-菜单三级权限，动态菜单+路由守卫',
    '✓ 文件管理：文件上传/删除/MD5校验',
], size=12, color=MID_TEXT)
add_text(s, Inches(7.0), Inches(1.5), Inches(5.5), Inches(0.45), '未来展望', size=20, color=PRIMARY_DARK, bold=True)
add_multitext(s, Inches(7.0), Inches(2.1), Inches(5.5), Inches(3.8), [
    '1. 支持更多数据格式（Shapefile、GeoJSON等）',
    '2. 增加作业质量评估与统计分析可视化模块',
    '3. 引入机器学习进行轨迹异常检测',
    '4. 支持多地块作业数据对比分析',
    '5. 移动端适配与PWA渐进式应用',
    '6. 接入实时GPS数据流，实现在线作业监控',
    '7. 优化大数据量渲染性能（虚拟列表、WebGL）',
    '8. 导出轨迹数据为PDF报告',
    '9. 农田边界识别与面积自动计算',
    '10. 作业路径规划与优化建议',
], size=12, color=MID_TEXT)
add_rect(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.7), PRIMARY)
add_text(s, Inches(0.9), Inches(6.18), Inches(11.5), Inches(0.55),
         '技术栈：Python Django 4.2 + Vue 3 + Vite + Element Plus + Leaflet + SQLite + Django REST Framework',
         size=13, color=WHITE, align=PP_ALIGN.CENTER)
page_num(s, 21, TOTAL)

# =====================================================================
# SLIDE 22 – 感谢聆听
# =====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s, DARK_BG)
add_rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PRIMARY)
add_text(s, Inches(1.5), Inches(2.0), Inches(10.5), Inches(1.0),
         '感谢聆听', size=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(3.1), Inches(10.5), Inches(0.7),
         'THANK YOU', size=22, color=RGBColor(0x90,0xAA,0xCC), align=PP_ALIGN.CENTER)
add_rect(s, Inches(5.8), Inches(4.0), Inches(1.733), Inches(0.04), PRIMARY)
add_text(s, Inches(1.5), Inches(4.5), Inches(10.5), Inches(0.6),
         '农机作业数据处理与分析软件', size=22, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1.5), Inches(5.1), Inches(10.5), Inches(0.6),
         '首都师范大学  万爱华', size=18, color=RGBColor(0xCC,0xDD,0xFF), align=PP_ALIGN.CENTER)
add_rect(s, Inches(0), Inches(7.2), Inches(13.333), Inches(0.3), PRIMARY)
page_num(s, 22, TOTAL)

# ====== SAVE ======
prs.save(OUTPUT_PATH)
print(f'Done! {len(prs.slides)} slides → {OUTPUT_PATH}')
